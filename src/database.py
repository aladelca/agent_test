import sqlite3
import json
import os
from datetime import datetime
from typing import List, Dict, Optional
from langchain.chat_models import ChatOpenAI
from langchain.schema import HumanMessage
from langchain.schema.messages import BaseMessage

class CourseDatabase:
    def __init__(self, db_path="courses.db"):
        self.db_path = db_path
        self.llm = ChatOpenAI(model_name="gpt-3.5-turbo", temperature=0)
        self.init_db()

    def get_connection(self):
        return sqlite3.connect(self.db_path)

    def init_db(self):
        """Initializes the database with the required tables"""
        with self.get_connection() as conn:
            cursor = conn.cursor()
            
            # Courses table
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS courses (
                    id INTEGER PRIMARY KEY,
                    name TEXT NOT NULL,
                    section TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            """)
            
            # Documents table
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS documents (
                    id INTEGER PRIMARY KEY,
                    course_id INTEGER,
                    title TEXT NOT NULL,
                    file_path TEXT NOT NULL,
                    category TEXT NOT NULL,
                    ciclo TEXT,
                    modulo TEXT,
                    section TEXT,
                    uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    FOREIGN KEY (course_id) REFERENCES courses(id)
                )
            """)
            
            # Updates table
            cursor.execute("""
                CREATE TABLE IF NOT EXISTS updates (
                    id INTEGER PRIMARY KEY,
                    course_id INTEGER,
                    content TEXT NOT NULL,
                    category TEXT NOT NULL,
                    ciclo TEXT,
                    modulo TEXT,
                    section TEXT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    FOREIGN KEY (course_id) REFERENCES courses(id)
                )
            """)
            
            conn.commit()

    def get_courses(self) -> Dict:
        """Gets all courses and their updates"""
        with self.get_connection() as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT id, name, section FROM courses")
            courses = {}
            
            for course_id, name, course_section in cursor.fetchall():
                # Get updates
                cursor.execute("""
                    SELECT content, category, ciclo, modulo, section, created_at 
                    FROM updates 
                    WHERE course_id = ?
                    ORDER BY created_at DESC
                """, (course_id,))
                
                updates = []
                for content, category, ciclo, modulo, update_section, created_at in cursor.fetchall():
                    updates.append({
                        'content': content,
                        'category': category,
                        'ciclo': ciclo,
                        'modulo': modulo,
                        'section': course_section,  # Use course section
                        'timestamp': created_at
                    })
                
                # Get documents
                cursor.execute("""
                    SELECT title, file_path, category, ciclo, modulo, section, uploaded_at 
                    FROM documents 
                    WHERE course_id = ?
                    ORDER BY uploaded_at DESC
                """, (course_id,))
                
                documents = []
                for title, file_path, category, ciclo, modulo, section, uploaded_at in cursor.fetchall():
                    documents.append({
                        'title': title,
                        'file_path': file_path,
                        'category': category,
                        'ciclo': ciclo,
                        'modulo': modulo,
                        'section': section,
                        'timestamp': uploaded_at
                    })
                
                courses[name] = {
                    'section': course_section,
                    'updates': updates,
                    'documents': documents,
                    'categories': list(set(update['category'] for update in updates)),
                    'last_update': updates[0]['timestamp'] if updates else None
                }
            
            return courses

    def add_update(self, course_name: str, section: str, content: str, category: str, ciclo: str, modulo: str) -> bool:
        """Adds an update to the course"""
        try:
            with self.get_connection() as conn:
                cursor = conn.cursor()
                
                # Get or create the course
                cursor.execute("SELECT id FROM courses WHERE name = ?", (course_name,))
                result = cursor.fetchone()
                
                if result:
                    course_id = result[0]
                    # Update section if needed
                    cursor.execute("UPDATE courses SET section = ? WHERE id = ?", (section, course_id))
                else:
                    cursor.execute("INSERT INTO courses (name, section) VALUES (?, ?)", (course_name, section))
                    course_id = cursor.lastrowid
                
                # Insert the update
                cursor.execute("""
                    INSERT INTO updates (course_id, content, category, ciclo, modulo, section)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (course_id, content, category, ciclo, modulo, section))
                
                conn.commit()
                print(f"Update saved successfully for course {course_name}")
                return True
        except Exception as e:
            print(f"Error adding update: {e}")
            return False

    def add_document(self, course_name: str, section: str, file_path: str, title: str, category: str, ciclo: str, modulo: str) -> bool:
        """Adds a document to the course"""
        try:
            with self.get_connection() as conn:
                cursor = conn.cursor()
                
                # Get or create the course
                cursor.execute("SELECT id FROM courses WHERE name = ?", (course_name,))
                result = cursor.fetchone()
                
                if result:
                    course_id = result[0]
                else:
                    cursor.execute("INSERT INTO courses (name, section) VALUES (?, ?)", (course_name, section))
                    course_id = cursor.lastrowid
                
                # Insert the document
                cursor.execute("""
                    INSERT INTO documents (course_id, title, file_path, category, ciclo, modulo)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (course_id, title, file_path, category, ciclo, modulo))
                
                conn.commit()
                return True
        except Exception as e:
            print(f"Error adding document: {e}")
            return False

    def detect_category(self, content: str, file_name: str = None) -> str:
        """
        Uses LLM to detect the content or document category.
        Returns the most appropriate category.
        """
        prompt = f"""Analiza el siguiente contenido y determina su categoría más apropiada.
        
        {'Nombre del archivo: ' + file_name if file_name else 'Contenido: ' + content}
        
        Categorías disponibles:
        - EVALUACIÓN: Exámenes, prácticas, evaluaciones
        - CLASE: Materiales de clase, presentaciones, guías
        - TAREA: Tareas, trabajos, proyectos
        - SÍLABO: Sílabos, cronogramas, objetivos
        - CRONOGRAMA: Fechas importantes, calendarios
        - GENERAL: Otros anuncios o información general
        
        Responde ÚNICAMENTE con el nombre de la categoría más apropiada.
        No incluyas explicaciones ni texto adicional.
        """
        
        response = self.llm.invoke([HumanMessage(content=prompt)])
        category = response.content.strip()
        
        # Validate that the category is valid
        valid_categories = ['EVALUACIÓN', 'CLASE', 'TAREA', 'SÍLABO', 'CRONOGRAMA', 'GENERAL']
        return category if category in valid_categories else 'GENERAL'

    def migrate_from_json_file(self, json_file: str):
        """Migrates data from a JSON file to the database"""
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            with self.get_connection() as conn:
                cursor = conn.cursor()
                
                for course_name, course_data in data.items():
                    # Insert the course
                    section = course_data.get('section', '')
                    cursor.execute("""
                        INSERT INTO courses (name, section)
                        VALUES (?, ?)
                    """, (course_name, section))
                    course_id = cursor.lastrowid
                    
                    # Insert updates
                    for update in course_data.get('updates', []):
                        cursor.execute("""
                            INSERT INTO updates (course_id, content, category, ciclo, modulo, section, created_at)
                            VALUES (?, ?, ?, ?, ?, ?, ?)
                        """, (
                            course_id,
                            update['content'],
                            update['category'],
                            update['ciclo'],
                            update['modulo'],
                            section,  # Use course section
                            update['timestamp']
                        ))
                
                conn.commit()
                print("Migration completed successfully")
        except Exception as e:
            print(f"Error during migration: {e}")

    def get_course_info(self, course_name: str, ciclo: Optional[str] = None, modulo: Optional[str] = None, section: Optional[str] = None) -> Optional[Dict]:
        """Gets all information for a course, including document contents"""
        try:
            with self.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute("SELECT id, section FROM courses WHERE name = ?", (course_name,))
                result = cursor.fetchone()
                
                if not result:
                    return None
                
                course_id, course_section = result
                
                # Obtener actualizaciones
                query = """
                    SELECT content, category, ciclo, modulo, section, created_at 
                    FROM updates 
                    WHERE course_id = ?
                """
                params = [course_id]
                
                if ciclo:
                    query += " AND ciclo = ?"
                    params.append(ciclo)
                if modulo:
                    query += " AND modulo = ?"
                    params.append(modulo)
                if section:
                    query += " AND section = ?"
                    params.append(section)
                
                query += " ORDER BY created_at DESC"
                cursor.execute(query, tuple(params))
                
                updates = []
                for content, category, update_ciclo, update_modulo, update_section, created_at in cursor.fetchall():
                    updates.append({
                        'content': content,
                        'category': category,
                        'ciclo': update_ciclo,
                        'modulo': update_modulo,
                        'section': update_section,
                        'timestamp': created_at
                    })
                
                # Obtener documentos
                query = """
                    SELECT title, file_path, category, ciclo, modulo, section, uploaded_at 
                    FROM documents 
                    WHERE course_id = ?
                """
                params = [course_id]
                
                if ciclo:
                    query += " AND ciclo = ?"
                    params.append(ciclo)
                if modulo:
                    query += " AND modulo = ?"
                    params.append(modulo)
                if section:
                    query += " AND section = ?"
                    params.append(section)
                
                query += " ORDER BY uploaded_at DESC"
                cursor.execute(query, tuple(params))
                
                documents = []
                for title, file_path, category, doc_ciclo, doc_modulo, doc_section, uploaded_at in cursor.fetchall():
                    documents.append({
                        'title': title,
                        'file_path': file_path,
                        'category': category,
                        'ciclo': doc_ciclo,
                        'modulo': doc_modulo,
                        'section': doc_section,
                        'timestamp': uploaded_at
                    })
                
                return {
                    'name': course_name,
                    'section': section or course_section,
                    'updates': updates,
                    'documents': documents,
                    'categories': list(set(update['category'] for update in updates)),
                    'last_update': updates[0]['timestamp'] if updates else None
                }
                
        except Exception as e:
            print(f"Error getting course info: {e}")
            return None

    def read_document_content(self, file_path: str) -> str:
        """Lee el contenido de un documento y lo devuelve como texto"""
        try:
            # Verificar que el archivo existe
            if not os.path.exists(file_path):
                return ""

            # Obtener la extensión del archivo
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()

            # Manejar diferentes tipos de archivos
            if ext in ['.txt', '.md', '.csv']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif ext == '.pdf':
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            elif ext in ['.doc', '.docx']:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            elif ext in ['.xls', '.xlsx']:
                import pandas as pd
                df = pd.read_excel(file_path)
                return df.to_string()
            elif ext in ['.ppt', '.pptx']:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                return f"[Tipo de archivo no soportado: {ext}]"
        except Exception as e:
            print(f"Error reading document content: {e}")
            return ""

    def get_document_content(self, course_name: str, document_title: str) -> Optional[str]:
        """Obtiene el contenido de un documento específico"""
        try:
            with self.get_connection() as conn:
                cursor = conn.cursor()
                
                # Obtener el ID del curso
                cursor.execute("SELECT id FROM courses WHERE name = ?", (course_name,))
                result = cursor.fetchone()
                
                if not result:
                    return None
                
                course_id = result[0]
                
                # Obtener la ruta del archivo
                cursor.execute("""
                    SELECT file_path 
                    FROM documents 
                    WHERE course_id = ? AND title = ?
                """, (course_id, document_title))
                
                result = cursor.fetchone()
                if not result:
                    return None
                
                file_path = result[0]
                return self.read_document_content(file_path)
        except Exception as e:
            print(f"Error getting document content: {e}")
            return None 